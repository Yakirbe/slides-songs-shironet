#!/usr/bin/env python3
"""
Lyrics PPTX Generator
---------------------
Creates a PowerPoint presentation from lyrics text files.
Automatically uses multi-column layout when lyrics are too long.
Uses Unsplash for thematic background images.

Usage: python create_lyrics_pptx.py
"""

import os
import hashlib
import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# Cache directory for downloaded images
CACHE_DIR = Path(__file__).parent / 'image_cache'


def get_search_terms(title: str) -> str:
    """Extract meaningful search terms from song title."""
    # Remove artist names after dash/hyphen
    if ' - ' in title:
        title = title.split(' - ')[0]
    if ' / ' in title:
        title = title.split(' / ')[0]
    
    # Common Hebrew to English theme mappings for better Unsplash results
    theme_map = {
        'אהבה': 'love romance',
        'לב': 'heart love',
        'שמש': 'sun sunshine',
        'ירח': 'moon night',
        'ים': 'sea ocean',
        'גשם': 'rain',
        'לילה': 'night stars',
        'בית': 'home house',
        'דרך': 'road journey path',
        'חלום': 'dream dreaming',
        'שיר': 'music singing',
        'ריקוד': 'dance dancing',
        'אביב': 'spring flowers',
        'סתיו': 'autumn fall leaves',
        'חורף': 'winter snow',
        'קיץ': 'summer beach',
    }
    
    # Check if title contains Hebrew theme words
    for hebrew, english in theme_map.items():
        if hebrew in title:
            return english
    
    # For English titles, use the title itself
    if not is_hebrew(title):
        return title
    
    # Default fallback for Hebrew songs
    return 'music abstract'


def download_background_image(song_index: int, width: int = 1920, height: int = 1080) -> Path | None:
    """Download a background image from Lorem Picsum (free, high quality)."""
    CACHE_DIR.mkdir(exist_ok=True)
    
    # Use song index as seed for consistent image per song
    cache_path = CACHE_DIR / f"bg_{song_index:03d}.jpg"
    
    # Return cached image if exists
    if cache_path.exists():
        return cache_path
    
    try:
        # Lorem Picsum with seed for reproducible images
        url = f"https://picsum.photos/seed/{song_index}/{width}/{height}"
        response = requests.get(url, timeout=20, allow_redirects=True)
        
        if response.status_code == 200 and 'image' in response.headers.get('content-type', ''):
            cache_path.write_bytes(response.content)
            return cache_path
    except Exception as e:
        print(f"  Warning: Could not download image: {e}")
    
    return None


def get_optimal_layout(lines: list[str]) -> dict:
    """Determine optimal column layout based on line count."""
    line_count = len(lines)
    
    if line_count <= 15:
        return {'columns': 1, 'font_size': 22}
    elif line_count <= 28:
        return {'columns': 2, 'font_size': 18}
    elif line_count <= 42:
        return {'columns': 3, 'font_size': 16}
    elif line_count <= 56:
        return {'columns': 4, 'font_size': 14}
    elif line_count <= 72:
        return {'columns': 4, 'font_size': 12}
    else:
        return {'columns': 5, 'font_size': 10}


def split_into_columns(lines: list[str], num_columns: int) -> list[list[str]]:
    """Split lines evenly into columns."""
    lines_per_col = (len(lines) + num_columns - 1) // num_columns
    columns = []
    for i in range(num_columns):
        start = i * lines_per_col
        end = start + lines_per_col
        columns.append(lines[start:end])
    return columns


def is_hebrew_song(title: str, lyrics_lines: list[str]) -> bool:
    """Check if the song is primarily Hebrew based on title and content."""
    # Check title first
    if is_hebrew(title):
        return True
    # Check first few lines of lyrics
    hebrew_lines = sum(1 for line in lyrics_lines[:5] if is_hebrew(line))
    return hebrew_lines >= 2


def add_lyrics_slide(prs: Presentation, title: str, lyrics_lines: list[str], song_index: int = 0):
    """Add a slide with lyrics, using multi-column layout if needed."""
    # Use blank layout
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Determine if RTL (Hebrew song)
    is_rtl = is_hebrew_song(title, lyrics_lines)
    
    # Try to add background image from Lorem Picsum
    bg_image = download_background_image(song_index)
    
    # Add dark background first (always)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, slide_width, slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 25, 35)
    background.line.fill.background()
    
    # Add background image if available (on top of dark bg)
    has_bg_image = False
    if bg_image:
        try:
            slide.shapes.add_picture(
                str(bg_image), 0, 0, slide_width, slide_height
            )
            has_bg_image = True
        except Exception as e:
            print(f"  Warning: Could not add image: {e}")
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2),
        slide_width - Inches(0.6), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(40, 40, 60)  # Dark blue-gray
    title_para.alignment = PP_ALIGN.CENTER
    
    # Get layout config
    layout = get_optimal_layout(lyrics_lines)
    num_columns = layout['columns']
    font_size = layout['font_size']
    
    # Split into sections and assign colors
    sections = split_into_sections(lyrics_lines)
    colored_lines = assign_section_colors(sections)
    
    # Split into columns (now with color info)
    columns = split_into_columns(colored_lines, num_columns)
    
    # Calculate column dimensions
    margin = Inches(0.3)
    content_top = Inches(0.9)
    content_height = slide_height - content_top - Inches(0.3)
    available_width = slide_width - (2 * margin)
    col_gap = Inches(0.2)
    col_width = (available_width - (col_gap * (num_columns - 1))) / num_columns
    
    # Add text boxes for each column
    for i, col_lines in enumerate(columns):
        # For RTL: first column (i=0) goes on RIGHT side, then moves left
        # For LTR: first column (i=0) goes on LEFT side, then moves right
        if is_rtl:
            left = slide_width - margin - col_width - (i * (col_width + col_gap))
        else:
            left = margin + (i * (col_width + col_gap))
        
        # Add semi-transparent background box for text readability
        if has_bg_image:
            text_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.1), content_top - Inches(0.05),
                col_width + Inches(0.2), content_height + Inches(0.1)
            )
            text_bg.fill.solid()
            text_bg.fill.fore_color.rgb = RGBColor(30, 30, 50)  # Dark blue-gray, softer than black
            text_bg.line.fill.background()
            # Set corner radius
            text_bg.adjustments[0] = 0.08
        
        text_box = slide.shapes.add_textbox(
            left, content_top,
            col_width, content_height
        )
        
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        # Add lyrics text with section colors
        for j, item in enumerate(col_lines):
            line, color = item  # Each item is (line_text, color)
            
            if j == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            
            para.text = line
            para.font.size = Pt(font_size)
            para.font.color.rgb = color  # Use section color
            para.alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
            para.space_after = Pt(4)


def is_hebrew(text: str) -> bool:
    """Check if text contains Hebrew characters."""
    for char in text:
        if '\u0590' <= char <= '\u05FF':
            return True
    return False


def clean_punctuation(text: str) -> str:
    """Remove commas and periods from text."""
    return text.replace(',', '').replace('.', '')


# Color palette for alternating verse/chorus sections
SECTION_COLORS = [
    RGBColor(240, 240, 255),   # White-blue (default)
    RGBColor(255, 220, 180),   # Warm peach
    RGBColor(180, 255, 200),   # Mint green
    RGBColor(255, 200, 220),   # Soft pink
    RGBColor(200, 220, 255),   # Light blue
    RGBColor(255, 255, 180),   # Soft yellow
]


def split_into_sections(lines: list[str]) -> list[list[str]]:
    """Split lyrics into sections (verses/choruses) based on empty lines."""
    sections = []
    current_section = []
    
    for line in lines:
        if not line.strip():
            if current_section:
                sections.append(current_section)
                current_section = []
        else:
            current_section.append(line)
    
    if current_section:
        sections.append(current_section)
    
    return sections


def assign_section_colors(sections: list[list[str]]) -> list[tuple[str, RGBColor]]:
    """Assign alternating colors to each line based on its section."""
    result = []
    for i, section in enumerate(sections):
        color = SECTION_COLORS[i % len(SECTION_COLORS)]
        for line in section:
            result.append((line, color))
        # Add empty line between sections (except after last)
        if i < len(sections) - 1:
            result.append(('', SECTION_COLORS[0]))  # Empty line
    return result


def parse_lyrics_file(filepath: Path) -> tuple[str, list[str]]:
    """Parse a lyrics file and return title and lines (preserving empty lines for verse divisions)."""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    lines = content.strip().split('\n')
    
    # Extract title (first line starting with #)
    title = filepath.stem.replace('_', ' ')
    lyrics_lines = []
    found_title = False
    
    for line in lines:
        if line.startswith('# '):
            title = line[2:].strip()
            found_title = True
        else:
            # Keep all lines including empty ones (preserve verse/chorus structure)
            # Remove commas and periods
            lyrics_lines.append(clean_punctuation(line))
    
    # Remove leading/trailing empty lines only
    while lyrics_lines and not lyrics_lines[0].strip():
        lyrics_lines.pop(0)
    while lyrics_lines and not lyrics_lines[-1].strip():
        lyrics_lines.pop()
    
    return title, lyrics_lines


def main():
    base_dir = Path(__file__).parent
    lyrics_dir = base_dir / 'lyrics'
    output_path = base_dir / 'output' / 'songs_presentation.pptx'
    
    # Ensure output directory exists
    output_path.parent.mkdir(exist_ok=True)
    
    # Create presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Get all lyrics files
    lyrics_files = sorted(lyrics_dir.glob('*.txt'))
    print(f"Found {len(lyrics_files)} lyrics files")
    
    # Add title slide
    slide_layout = prs.slide_layouts[6]  # Blank
    title_slide = prs.slides.add_slide(slide_layout)
    
    # Title slide background
    bg = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 20, 35)
    bg.line.fill.background()
    
    # Title text
    title_box = title_slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        prs.slide_width - Inches(1), Inches(2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎵 Song Lyrics Collection"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 220, 120)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = f"{len(lyrics_files)} Songs"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(180, 180, 200)
    p2.alignment = PP_ALIGN.CENTER
    
    # Add a slide for each song
    for idx, lyrics_file in enumerate(lyrics_files):
        print(f"Adding: {lyrics_file.stem}")
        title, lyrics_lines = parse_lyrics_file(lyrics_file)
        
        if lyrics_lines:
            add_lyrics_slide(prs, title, lyrics_lines, song_index=idx)
    
    # Save presentation
    prs.save(str(output_path))
    print(f"\n✓ Saved to: {output_path}")


if __name__ == '__main__':
    main()
