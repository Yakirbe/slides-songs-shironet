"""
Lyrics Slide Generator - Backend API
FastAPI server for searching lyrics and generating PPTX slides.
"""

import os
import re
import hashlib
import tempfile
from pathlib import Path
from typing import Optional

import requests
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


app = FastAPI(title="Lyrics Slide Generator")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Cache directory
CACHE_DIR = Path(__file__).parent / "cache"
CACHE_DIR.mkdir(exist_ok=True)


class SearchRequest(BaseModel):
    query: str


class GenerateRequest(BaseModel):
    title: str
    url: str


class SearchResult(BaseModel):
    title: str
    artist: str
    url: str


# ============== Shironet Functions ==============

def search_shironet(query: str) -> list[SearchResult]:
    """Search Shironet for songs matching query."""
    results = []
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=False)
        context = browser.new_context(
            user_agent='Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36'
        )
        page = context.new_page()
        
        try:
            search_url = f"https://shironet.mako.co.il/search?q={query}"
            page.goto(search_url, wait_until='domcontentloaded', timeout=30000)
            page.wait_for_timeout(2000)
            
            # Find song links
            links = page.query_selector_all('a[href*="type=lyrics"]')
            
            for link in links[:10]:  # Limit to 10 results
                href = link.get_attribute('href')
                text = link.inner_text().strip()
                
                if href and text and 'wrkid' in href:
                    # Parse title and artist from text
                    if ' - ' in text:
                        parts = text.split(' - ', 1)
                        title = parts[0].strip()
                        artist = parts[1].strip() if len(parts) > 1 else ''
                    else:
                        title = text
                        artist = ''
                    
                    full_url = f"https://shironet.mako.co.il{href}" if href.startswith('/') else href
                    results.append(SearchResult(title=title, artist=artist, url=full_url))
                    
        except Exception as e:
            print(f"Search error: {e}")
        finally:
            browser.close()
    
    return results


def extract_lyrics(url: str) -> tuple[str, list[str]]:
    """Extract lyrics from a Shironet URL."""
    title = ""
    lyrics_lines = []
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=False)
        context = browser.new_context(
            user_agent='Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36'
        )
        page = context.new_page()
        
        try:
            page.goto(url, wait_until='domcontentloaded', timeout=30000)
            page.wait_for_timeout(2000)
            
            # Get title
            title_elem = page.query_selector('h1, .artist_song_name_txt')
            if title_elem:
                title = title_elem.inner_text().strip()
            
            # Get lyrics
            lyrics_elem = page.query_selector('span.artist_lyrics_text')
            if lyrics_elem:
                lyrics_text = lyrics_elem.inner_text()
                lyrics_lines = lyrics_text.split('\n')
                
        except Exception as e:
            print(f"Extract error: {e}")
        finally:
            browser.close()
    
    return title, lyrics_lines


# ============== PPTX Functions ==============

SECTION_COLORS = [
    RGBColor(240, 240, 255),
    RGBColor(255, 220, 180),
    RGBColor(180, 255, 200),
    RGBColor(255, 200, 220),
    RGBColor(200, 220, 255),
    RGBColor(255, 255, 180),
]


def is_hebrew(text: str) -> bool:
    for char in text:
        if '\u0590' <= char <= '\u05FF':
            return True
    return False


def clean_punctuation(text: str) -> str:
    return text.replace(',', '').replace('.', '')


def get_optimal_layout(lines: list[str]) -> dict:
    line_count = len(lines)
    if line_count <= 15:
        return {'columns': 1, 'font_size': 22}
    elif line_count <= 28:
        return {'columns': 2, 'font_size': 18}
    elif line_count <= 42:
        return {'columns': 3, 'font_size': 16}
    elif line_count <= 56:
        return {'columns': 4, 'font_size': 14}
    else:
        return {'columns': 5, 'font_size': 10}


def split_into_sections(lines: list[str]) -> list[list[str]]:
    sections = []
    current_section = []
    for line in lines:
        if not line.strip():
            if current_section:
                sections.append(current_section)
                current_section = []
        else:
            current_section.append(clean_punctuation(line))
    if current_section:
        sections.append(current_section)
    return sections


def download_background_image(seed: int) -> Optional[Path]:
    cache_path = CACHE_DIR / f"bg_{seed:03d}.jpg"
    if cache_path.exists():
        return cache_path
    try:
        url = f"https://picsum.photos/seed/{seed}/1920/1080"
        response = requests.get(url, timeout=20, allow_redirects=True)
        if response.status_code == 200:
            cache_path.write_bytes(response.content)
            return cache_path
    except Exception:
        pass
    return None


def generate_pptx(title: str, lyrics_lines: list[str]) -> Path:
    """Generate a PPTX slide for the given lyrics."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    is_rtl = is_hebrew(title)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(25, 25, 35)
    bg.line.fill.background()
    
    # Background image
    bg_image = download_background_image(hash(title) % 1000)
    has_bg_image = False
    if bg_image:
        try:
            slide.shapes.add_picture(str(bg_image), 0, 0, slide_width, slide_height)
            has_bg_image = True
        except Exception:
            pass
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), slide_width - Inches(0.6), Inches(0.6))
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(40, 40, 60)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Sections and colors
    sections = split_into_sections(lyrics_lines)
    colored_lines = []
    for i, section in enumerate(sections):
        color = SECTION_COLORS[i % len(SECTION_COLORS)]
        for line in section:
            colored_lines.append((line, color))
        if i < len(sections) - 1:
            colored_lines.append(('', SECTION_COLORS[0]))
    
    # Layout
    layout = get_optimal_layout(colored_lines)
    num_columns = layout['columns']
    font_size = layout['font_size']
    
    # Split into columns
    lines_per_col = (len(colored_lines) + num_columns - 1) // num_columns
    columns = []
    for i in range(num_columns):
        start = i * lines_per_col
        end = start + lines_per_col
        columns.append(colored_lines[start:end])
    
    margin = Inches(0.3)
    content_top = Inches(0.9)
    content_height = slide_height - content_top - Inches(0.3)
    available_width = slide_width - (2 * margin)
    col_gap = Inches(0.2)
    col_width = (available_width - (col_gap * (num_columns - 1))) / num_columns
    
    for i, col_lines in enumerate(columns):
        if is_rtl:
            left = slide_width - margin - col_width - (i * (col_width + col_gap))
        else:
            left = margin + (i * (col_width + col_gap))
        
        if has_bg_image:
            text_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.1), content_top - Inches(0.05),
                col_width + Inches(0.2), content_height + Inches(0.1)
            )
            text_bg.fill.solid()
            text_bg.fill.fore_color.rgb = RGBColor(30, 30, 50)
            text_bg.line.fill.background()
            text_bg.adjustments[0] = 0.08
        
        text_box = slide.shapes.add_textbox(left, content_top, col_width, content_height)
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for j, item in enumerate(col_lines):
            line, color = item
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(font_size)
            para.font.color.rgb = color
            para.alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
            para.space_after = Pt(4)
    
    # Save to temp file
    output_path = Path(tempfile.mktemp(suffix='.pptx'))
    prs.save(str(output_path))
    return output_path


# ============== API Endpoints ==============

@app.get("/")
def root():
    return {"message": "Lyrics Slide Generator API", "version": "1.0"}


@app.post("/search")
def search(request: SearchRequest):
    """Search for songs on Shironet."""
    results = search_shironet(request.query)
    return {"results": results}


@app.post("/generate")
def generate(request: GenerateRequest):
    """Generate PPTX slide from a Shironet URL."""
    title, lyrics = extract_lyrics(request.url)
    
    if not lyrics:
        raise HTTPException(status_code=404, detail="Could not extract lyrics")
    
    # Use provided title if extraction failed
    if not title:
        title = request.title
    
    pptx_path = generate_pptx(title, lyrics)
    
    return FileResponse(
        pptx_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{title.replace(' ', '_')}.pptx"
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
